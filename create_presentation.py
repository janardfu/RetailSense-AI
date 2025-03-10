from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os
import re

# Define standard margins and layout settings
MARGIN_LEFT = Inches(0.5)
MARGIN_RIGHT = Inches(0.5)
MARGIN_TOP = Inches(0.5)
MARGIN_BOTTOM = Inches(0.5)
CONTENT_WIDTH = Inches(12.33)  # 13.33 - (0.5 * 2) for margins
CONTENT_HEIGHT = Inches(6.5)   # 7.5 - (0.5 * 2) for margins

def read_presentation_content(file_path='RetailSense_Presentation.txt'):
    """Read and parse presentation content from text file"""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split into sections based on markdown-style headers
    sections = {}
    current_section = None
    current_subsection = None
    current_content = []
    
    lines = content.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        
        # Skip empty lines
        if not line:
            i += 1
            continue
            
        # Main section header (with number and =)
        if i + 1 < len(lines) and '=' in lines[i + 1]:
            if current_section and current_subsection and current_content:
                if current_section not in sections:
                    sections[current_section] = {}
                sections[current_section][current_subsection] = current_content
            current_section = line.strip()
            current_subsection = "Main"
            current_content = []
            i += 2  # Skip the === line
            
        # Special sections (PRESENTATION FLOW, VISUAL ELEMENTS, etc.)
        elif line.strip() == "PRESENTATION FLOW:":
            if current_section and current_subsection and current_content:
                if current_section not in sections:
                    sections[current_section] = {}
                sections[current_section][current_subsection] = current_content
            current_section = "PRESENTATION FLOW"
            current_subsection = "Main"
            current_content = []
            i += 1
            
        # Sub-section header (with •)
        elif line.startswith('•'):
            if current_subsection and current_content:
                if current_section not in sections:
                    sections[current_section] = {}
                sections[current_section][current_subsection] = current_content
            current_subsection = line.strip('• ').strip()
            current_content = []
            i += 1
            
        # Content lines
        else:
            if line.strip():
                current_content.append(line)
            i += 1
    
    # Add the last section/subsection
    if current_section and current_subsection and current_content:
        if current_section not in sections:
            sections[current_section] = {}
        sections[current_section][current_subsection] = current_content
    
    return sections

def apply_slide_layout(shape):
    """Apply standard layout to a shape"""
    shape.left = MARGIN_LEFT
    shape.top = MARGIN_TOP
    shape.width = CONTENT_WIDTH
    
def create_title_slide(prs):
    """Create title slide"""
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    # Style title
    title.text = "RetailSense"
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(64)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title.top = Inches(2)
    title.left = MARGIN_LEFT
    title.width = CONTENT_WIDTH
    
    # Style subtitle
    subtitle.text = "AI-Powered Inventory Management"
    subtitle_frame = subtitle.text_frame
    subtitle_frame.paragraphs[0].font.size = Pt(36)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle.top = Inches(3.5)
    subtitle.left = MARGIN_LEFT
    subtitle.width = CONTENT_WIDTH

def create_section_title(prs, title, subtitle=""):
    """Create section title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    # Style title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_shape.top = MARGIN_TOP
    title_shape.left = MARGIN_LEFT
    title_shape.width = CONTENT_WIDTH
    
    # Style subtitle
    subtitle_shape.text = subtitle
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.paragraphs[0].font.size = Pt(32)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    subtitle_shape.top = Inches(2)
    subtitle_shape.left = MARGIN_LEFT
    subtitle_shape.width = CONTENT_WIDTH
    
    return slide

def create_content_slide(prs, title, content_list):
    """Create content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    # Style title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_shape.top = MARGIN_TOP
    title_shape.left = MARGIN_LEFT
    title_shape.width = CONTENT_WIDTH
    
    # Style content
    tf = body_shape.text_frame
    tf.word_wrap = True
    body_shape.top = Inches(1.5)
    body_shape.left = MARGIN_LEFT
    body_shape.width = CONTENT_WIDTH
    body_shape.height = CONTENT_HEIGHT - Inches(1.5)
    
    for idx, item in enumerate(content_list):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Handle bullet points and sub-bullet points
        item = item.strip()
        if item.startswith('  -'):  # Sub-bullet point
            p.text = item.strip('- ').strip()
            p.level = 1
            p.alignment = PP_ALIGN.LEFT
        elif item.startswith('-'):  # Main bullet point
            p.text = item.strip('- ').strip()
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
        elif item.startswith('*'):  # Alternative bullet point
            p.text = item.strip('* ').strip()
            p.level = 1
            p.alignment = PP_ALIGN.LEFT
        else:  # Regular text
            p.text = item
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
        
        # Style paragraph
        p.font.size = Pt(24 if p.level == 0 else 20)
        p.space_before = Pt(12)
        p.space_after = Pt(6)
        
        # Add bullet points
        if p.level > 0:
            p.font.size = Pt(20)
            p.space_before = Pt(6)
            p.space_after = Pt(6)

def parse_presentation_flow(content):
    """Parse presentation flow section to get timing and structure"""
    flow = {}
    current_section = None
    
    for line in content:
        # Skip empty lines and section headers
        if not line.strip() or line.endswith(':'):
            continue
            
        # Match lines like "1. Introduction (2 min)"
        match = re.match(r'^\d+\.\s+(.*?)\s*\((\d+)\s*min\)', line)
        if match:
            section_name, timing = match.groups()
            current_section = section_name.strip()
            flow[current_section] = {
                'time': int(timing),
                'subsections': []
            }
        # Match subsection lines starting with "-"
        elif line.strip().startswith('-'):
            if current_section:
                subsection = line.strip('- ').strip()
                if ':' in subsection:  # Handle subsections with descriptions
                    subsection = subsection.split(':')[0].strip()
                flow[current_section]['subsections'].append(subsection)
    
    return flow

def find_section_content(sections, section_name, subsection_name):
    """Find content for a specific section and subsection"""
    # Map section names to their numbered versions
    section_mapping = {
        'Introduction': '1. INNOVATION & CREATIVITY',
        'Innovation Showcase': '1. INNOVATION & CREATIVITY',
        'Technical Demo': '2. AI MODEL USAGE',
        'Business Impact': '4. BUSINESS USE CASE',
        'Q&A': '5. DEMO & PRESENTATION'
    }
    
    # Map subsection names to their actual titles
    subsection_mapping = {
        'Problem Statement': 'Current market challenges',
        'Solution Overview': 'RetailSense platform',
        'Unique Features': 'Unique Value Proposition',
        'Technical Innovation': 'Advanced AI Implementation',
        'Dashboard': 'Live Dashboard Demo',
        'AI Predictions': 'AI-Powered Analytics',
        'ROI Metrics': 'ROI Potential',
        'Market Potential': 'Market Application',
        'Technical Details': 'Architecture',
        'Implementation': 'Code Quality'
    }
    
    # Get the mapped section name
    mapped_section = section_mapping.get(section_name, section_name)
    # Get the mapped subsection name
    mapped_subsection = subsection_mapping.get(subsection_name, subsection_name)
    
    # First try exact match with mapped names
    for section_title, section_data in sections.items():
        if mapped_section in section_title and isinstance(section_data, dict):
            for subsec_title, content in section_data.items():
                if mapped_subsection.lower() in subsec_title.lower():
                    return content
    
    # If not found, try fuzzy match with original names
    for section_title, section_data in sections.items():
        if isinstance(section_data, dict):
            for subsec_title, content in section_data.items():
                if subsection_name.lower() in subsec_title.lower():
                    return content
                elif any(word.lower() in subsec_title.lower() for word in subsection_name.split()):
                    return content
    
    return None

def save_presentation(prs, base_filename='RetailSense_Presentation'):
    """Save presentation with proper file path logging"""
    # Get absolute path of current directory
    current_dir = os.path.abspath(os.path.dirname(__file__))
    
    # Create Presentations folder if it doesn't exist
    presentations_dir = os.path.join(current_dir, 'Presentations')
    if not os.path.exists(presentations_dir):
        os.makedirs(presentations_dir)
        print(f"\nCreated Presentations directory at: {presentations_dir}")
    
    # Create filename with timestamp
    from datetime import datetime
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{base_filename}_{timestamp}.pptx"
    filepath = os.path.join(presentations_dir, filename)
    
    try:
        # Try to save in Presentations directory
        prs.save(filepath)
        print(f"\nPresentation saved successfully!")
        print(f"File location: {filepath}")
        
        # List all presentations in the directory
        print("\nAvailable presentations in directory:")
        presentations = sorted([f for f in os.listdir(presentations_dir) if f.endswith('.pptx')])
        for idx, pres in enumerate(presentations, 1):
            pres_path = os.path.join(presentations_dir, pres)
            pres_size = os.path.getsize(pres_path) / 1024  # Convert to KB
            pres_time = datetime.fromtimestamp(os.path.getmtime(pres_path)).strftime("%Y-%m-%d %H:%M:%S")
            print(f"{idx}. {pres}")
            print(f"   Size: {pres_size:.2f} KB")
            print(f"   Modified: {pres_time}")
        
        return filepath
    except PermissionError:
        # If permission error, save to temp directory
        import tempfile
        temp_dir = tempfile.gettempdir()
        temp_filepath = os.path.join(temp_dir, filename)
        prs.save(temp_filepath)
        print(f"\nPermission denied in Presentations directory.")
        print(f"Presentation saved to temporary location: {temp_filepath}")
        return temp_filepath

def main():
    # Read content from file
    sections = read_presentation_content()
    
    # Debug output
    print("\nParsed sections:")
    for section_name, section_data in sections.items():
        print(f"\nSection: {section_name}")
        if isinstance(section_data, dict):
            for subsec_name, content in section_data.items():
                print(f"  Subsection: {subsec_name}")
                print(f"    Content lines: {len(content)}")
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions for 16:9 aspect ratio
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    create_title_slide(prs)
    print("\nCreated title slide")
    
    # Parse presentation flow
    flow = parse_presentation_flow(sections.get('PRESENTATION FLOW', {}).get('Main', []))
    
    # Debug output
    print("\nPresentation flow:")
    for section_name, section_info in flow.items():
        print(f"\nSection: {section_name}")
        print(f"  Time: {section_info['time']} minutes")
        print(f"  Subsections: {section_info['subsections']}")
    
    # Create slides based on presentation flow
    for section_name, section_info in flow.items():
        print(f"\nCreating section: {section_name}")
        
        # Create section title slide
        create_section_title(prs, section_name, f"{section_info['time']} minutes")
        print(f"  Created section title slide")
        
        # Create content slides for each subsection
        for subsection in section_info['subsections']:
            print(f"  Processing subsection: {subsection}")
            
            # Find matching content
            content = find_section_content(sections, section_name, subsection)
            
            if content:
                print(f"    Found content ({len(content)} lines)")
                create_content_slide(prs, subsection, content)
                print(f"    Created content slide")
            else:
                print(f"    No content found")
    
    # Add Visual Elements section
    if 'VISUAL ELEMENTS' in sections:
        print("\nProcessing Visual Elements section")
        create_section_title(prs, "Visual Elements", "Supporting Materials")
        
        for subsection_title, content in sections['VISUAL ELEMENTS'].items():
            if content and subsection_title != "Main":
                print(f"  Creating slide for: {subsection_title}")
                create_content_slide(prs, subsection_title, content)
    
    # Add Key Talking Points section
    if 'KEY TALKING POINTS' in sections:
        print("\nProcessing Key Talking Points section")
        create_section_title(prs, "Key Talking Points", "Important Discussion Areas")
        
        for subsection_title, content in sections['KEY TALKING POINTS'].items():
            if content and subsection_title != "Main":
                print(f"  Creating slide for: {subsection_title}")
                create_content_slide(prs, subsection_title, content)
    
    # Print total number of slides
    print(f"\nTotal slides created: {len(prs.slides)}")
    
    # Save the presentation with timestamp
    saved_filepath = save_presentation(prs)
    
    # Print file size
    file_size = os.path.getsize(saved_filepath) / 1024  # Convert to KB
    print(f"File size: {file_size:.2f} KB")

if __name__ == "__main__":
    main() 